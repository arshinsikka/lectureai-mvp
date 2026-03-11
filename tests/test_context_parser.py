"""
Tests for app/services/context_parser.py
"""
import os
import textwrap
from pathlib import Path

import pytest

from app.services.context_parser import (
    parse_context_file,
    parse_context_files,
    MAX_CONTEXT_CHARS,
    SUPPORTED_EXTENSIONS,
)

SAMPLE_PDF = Path("test_data/Lecture 10 - Neural Networks on Sequential Data.pdf")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_txt(path: Path, content: str) -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def _make_docx(path: Path, paragraphs: list[str]) -> Path:
    from docx import Document
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))
    return path


def _make_pptx(path: Path, slides: list[dict]) -> Path:
    """slides = [{"title": "...", "body": "...", "notes": "..."}]"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # Title and Content
    for s in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = s.get("title", "")
        slide.placeholders[1].text = s.get("body", "")
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# .txt
# ---------------------------------------------------------------------------

def test_parse_txt(tmp_path):
    f = _make_txt(tmp_path / "notes.txt", "Hello lecture notes.")
    text = parse_context_file(f)
    assert "Hello lecture notes." in text


# ---------------------------------------------------------------------------
# .docx
# ---------------------------------------------------------------------------

def test_parse_docx(tmp_path):
    f = _make_docx(tmp_path / "notes.docx", ["Introduction", "Main content here."])
    text = parse_context_file(f)
    assert "Introduction" in text
    assert "Main content here." in text


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------

def test_parse_pptx(tmp_path):
    f = _make_pptx(tmp_path / "slides.pptx", [
        {"title": "Slide One", "body": "Body text", "notes": "Speaker note here"},
        {"title": "Slide Two", "body": "More content", "notes": ""},
    ])
    text = parse_context_file(f)
    assert "Slide One" in text
    assert "Body text" in text
    assert "Speaker note here" in text
    assert "Slide Two" in text
    assert "More content" in text


# ---------------------------------------------------------------------------
# Real PDF from test_data/
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not SAMPLE_PDF.exists(), reason="Sample PDF not in test_data/")
def test_parse_real_pdf():
    text = parse_context_file(SAMPLE_PDF)
    assert len(text) > 100, "Expected substantial text from the lecture PDF"
    # Print preview so the tester can visually verify
    print(f"\n=== PDF extract preview (first 500 chars) ===\n{text[:500]}")
    print(f"Total chars extracted: {len(text)}")


# ---------------------------------------------------------------------------
# Unsupported extension
# ---------------------------------------------------------------------------

def test_parse_unsupported_extension(tmp_path):
    bad = tmp_path / "file.xyz"
    bad.write_bytes(b"data")
    with pytest.raises(ValueError, match="Unsupported context file format"):
        parse_context_file(bad)


# ---------------------------------------------------------------------------
# Corrupted file is skipped during parse_context_files
# ---------------------------------------------------------------------------

def test_corrupted_file_skipped(tmp_path, monkeypatch):
    import app.config as cfg_mod
    cfg_mod.get_settings.cache_clear()
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    cfg_mod.get_settings.cache_clear()

    session_id = "corrupt-test"
    ctx_dir = tmp_path / session_id / "context"
    ctx_dir.mkdir(parents=True)

    # valid file
    _make_txt(ctx_dir / "good.txt", "Good content here.")
    # corrupted PDF
    (ctx_dir / "bad.pdf").write_bytes(b"not a real pdf")

    result = parse_context_files(session_id)

    assert "Good content here." in result
    # should not raise — corrupt file is skipped

    cfg_mod.get_settings.cache_clear()
    monkeypatch.delenv("DATA_DIR", raising=False)


# ---------------------------------------------------------------------------
# No context directory → empty string
# ---------------------------------------------------------------------------

def test_no_context_dir(tmp_path, monkeypatch):
    import app.config as cfg_mod
    cfg_mod.get_settings.cache_clear()
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    cfg_mod.get_settings.cache_clear()

    result = parse_context_files("session-no-ctx")
    assert result == ""

    cfg_mod.get_settings.cache_clear()
    monkeypatch.delenv("DATA_DIR", raising=False)


# ---------------------------------------------------------------------------
# Truncation at MAX_CONTEXT_CHARS
# ---------------------------------------------------------------------------

def test_truncation(tmp_path, monkeypatch):
    import app.config as cfg_mod
    cfg_mod.get_settings.cache_clear()
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    cfg_mod.get_settings.cache_clear()

    session_id = "trunc-test"
    ctx_dir = tmp_path / session_id / "context"
    ctx_dir.mkdir(parents=True)
    # Write a file with way more than MAX_CONTEXT_CHARS
    big_text = "A" * (MAX_CONTEXT_CHARS + 10_000)
    _make_txt(ctx_dir / "big.txt", big_text)

    result = parse_context_files(session_id)
    assert len(result) <= MAX_CONTEXT_CHARS

    cfg_mod.get_settings.cache_clear()
    monkeypatch.delenv("DATA_DIR", raising=False)


# ---------------------------------------------------------------------------
# Combined output saved to context_text.txt
# ---------------------------------------------------------------------------

def test_output_file_written(tmp_path, monkeypatch):
    import app.config as cfg_mod
    cfg_mod.get_settings.cache_clear()
    monkeypatch.setenv("DATA_DIR", str(tmp_path))
    cfg_mod.get_settings.cache_clear()

    session_id = "save-test"
    ctx_dir = tmp_path / session_id / "context"
    ctx_dir.mkdir(parents=True)
    _make_txt(ctx_dir / "a.txt", "Alpha content.")
    _make_txt(ctx_dir / "b.txt", "Beta content.")

    parse_context_files(session_id)

    output = tmp_path / session_id / "context_text.txt"
    assert output.exists()
    saved = output.read_text(encoding="utf-8")
    assert "Alpha content." in saved
    assert "Beta content." in saved

    cfg_mod.get_settings.cache_clear()
    monkeypatch.delenv("DATA_DIR", raising=False)
