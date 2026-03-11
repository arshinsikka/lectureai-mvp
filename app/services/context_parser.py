import logging
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_CONTEXT_CHARS = 50_000

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt"}


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    pages: list[str] = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            pages.append(text)
    doc.close()

    if not pages:
        logger.warning("PDF '%s' yielded no extractable text (possibly scanned).", path.name)
        return ""

    return "\n\n".join(pages)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
        # Include speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")
        if parts:
            slides.append(f"--- Slide {slide_num} ---\n" + "\n".join(parts))

    return "\n\n".join(slides)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".txt": _extract_txt,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def parse_context_file(file_path: Path) -> str:
    """
    Extract text from a single context file.

    Returns extracted text, or empty string if the file has no usable content.
    Raises ValueError for unsupported extensions.
    """
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    if ext not in _EXTRACTORS:
        raise ValueError(
            f"Unsupported context file format '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    return _EXTRACTORS[ext](file_path)


def parse_context_files(session_id: str) -> str:
    """
    Extract and combine text from context files for a session.

    Searches two locations (in order, deduplicating by filename):
      1. data/{session_id}/context/   (files uploaded via the context/ subfolder)
      2. data/{session_id}/           (files placed directly in the session root)

    Saves the result to data/{session_id}/context_text.txt and returns the
    combined text (truncated to MAX_CONTEXT_CHARS).

    Returns an empty string when no context files are found (not an error).
    """
    from app.config import get_settings

    settings = get_settings()
    session_dir = settings.session_data_dir(session_id)
    context_subdir = session_dir / "context"

    # Collect candidate files from both locations, deduplicating by filename.
    # The context/ subfolder takes priority; session root is the fallback.
    seen_names: set[str] = set()
    candidate_files: list[Path] = []

    for search_dir in [context_subdir, session_dir]:
        if not search_dir.exists():
            continue
        for f in sorted(search_dir.iterdir()):
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS:
                if f.name not in seen_names:
                    seen_names.add(f.name)
                    candidate_files.append(f)

    if not candidate_files:
        logger.info("[%s] No context files found in session or context/ dir.", session_id)
        return ""

    chunks: list[str] = []
    for f in candidate_files:
        ext = f.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            logger.warning("[%s] Skipping unsupported context file: %s", session_id, f.name)
            continue
        try:
            text = parse_context_file(f)
        except Exception as exc:  # noqa: BLE001
            logger.error("[%s] Failed to parse '%s': %s — skipping.", session_id, f.name, exc)
            continue

        if not text:
            logger.warning("[%s] '%s' contained no extractable text — skipping.", session_id, f.name)
            continue

        logger.info("[%s] Extracted %d chars from '%s'.", session_id, len(text), f.name)
        chunks.append(f"=== {f.name} ===\n{text}")

    combined = "\n\n".join(chunks)

    if len(combined) > MAX_CONTEXT_CHARS:
        logger.warning(
            "[%s] Combined context truncated from %d → %d chars.",
            session_id, len(combined), MAX_CONTEXT_CHARS,
        )
        combined = combined[:MAX_CONTEXT_CHARS]

    # Save to session data dir
    output_path = settings.session_data_dir(session_id) / "context_text.txt"
    output_path.write_text(combined, encoding="utf-8")
    logger.info("[%s] Context saved → %s (%d chars)", session_id, output_path, len(combined))

    return combined
